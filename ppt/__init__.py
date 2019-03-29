from pptx import Presentation as ppt
from pptx.util import Inches
import os

DEFAULT_LEFT = Inches(2)
DEFAULT_TOP = Inches(1.5)
DEFAULT_HEIGHT = Inches(5.5)


class Presentation:
    def __init__(self, savepath, name):
        self.savepath = savepath
        self.name = name
        self.savename = os.path.join(savepath, name)
        self.presentation = ppt()

    def add_image_slide(
            self,
            image_path,
            title='',
            left=DEFAULT_LEFT,
            top=DEFAULT_TOP,
            height=DEFAULT_HEIGHT):
        slide_layout = self.presentation.slide_layouts[5]
        slide = self.presentation.slides.add_slide(slide_layout)
        slide.shapes.add_picture(image_path, left=left, top=top, height=height)
        slide_title = slide.shapes.title
        slide_title.text = title

    def save(self):
        self.presentation.save(self.savename)
